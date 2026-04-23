from pptx import Presentation

def criar_apresentacao():
    # Cria o objeto da apresentação
    prs = Presentation()

    # Função auxiliar para criar slides com tópicos (Bullet points)
    def adicionar_slide_topicos(titulo, topicos):
        layout = prs.slide_layouts[1] # Layout 1 é Título e Conteúdo
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = titulo
        
        corpo = slide.placeholders[1].text_frame
        for i, topico in enumerate(topicos):
            if i == 0:
                corpo.text = topico
            else:
                p = corpo.add_paragraph()
                p.text = topico

    # --- SLIDE 1: Capa ---
    layout_capa = prs.slide_layouts[0] # Layout 0 é Slide de Título
    slide_capa = prs.slides.add_slide(layout_capa)
    slide_capa.shapes.title.text = "Lucas Dias Calado"
    slide_capa.placeholders[1].text = "Trajetória Acadêmica, Docência e Pesquisa em Engenharia de Materiais\n[Coloque seu e-mail / LinkedIn aqui]"

    # --- SLIDE 2: Quem Sou Eu ---
    adicionar_slide_topicos(
        "Quem Sou Eu",
        [
            "Professor de Ensino Superior e Pesquisador.",
            "Forte background em Engenharia Mecânica e Ciência dos Materiais.",
            "Especialista em Manufatura Aditiva (Impressão 3D) de ligas metálicas.",
            "Atuação na interface entre a engenharia de materiais e o setor biomédico."
        ]
    )

    # --- SLIDE 3: Atuação Acadêmica e Ensino ---
    adicionar_slide_topicos(
        "Atuação Acadêmica e Ensino",
        [
            "Professor (Programa EPP): Foco nas áreas de Administração e Desenvolvimento de Sistemas.",
            "Tutor EAD: Suporte e orientação em programas de Engenharia e Arquitetura.",
            "Professor Auxiliar: Condução de aulas práticas em laboratórios de Engenharia Mecânica.",
            "Professor Técnico (FIEC): Ensino de disciplinas fundamentais como Desenho Técnico e Gestão da Qualidade."
        ]
    )

    # --- SLIDE 4: Pesquisa em Manufatura Aditiva ---
    adicionar_slide_topicos(
        "Pesquisa em Manufatura Aditiva",
        [
            "Estudo focado em ligas de titânio beta para aplicações biomédicas.",
            "Projeto em Andamento (24 meses): Desenvolvimento de implantes ósseos inovadores.",
            "Utilização da liga Ti-20Nb-6Ta impressa em 3D.",
            "Inovação: Aplicação de revestimento de cobre com propriedades bactericidas para maior segurança."
        ]
    )

    # --- SLIDE 5: Patente e Desenvolvimento Tecnológico ---
    adicionar_slide_topicos(
        "Patente e Desenvolvimento Tecnológico",
        [
            "Órgão: INPI (Instituto Nacional da Propriedade Industrial)",
            "Número de Registro: BR1020240200519 (Depósito em 27 de Setembro de 2024).",
            "Título: Método de controle de propriedades mecânicas em peças metálicas produzidas por manufatura aditiva de liga de titânio.",
            "Impacto: Avanço no controle de qualidade e resistência estrutural de peças impressas em 3D."
        ]
    )

    # --- SLIDE 6: Produção Científica ---
    adicionar_slide_topicos(
        "Publicações de Destaque",
        [
            "Artigo Recente: 'In situ modulation of mechanical properties in additively manufactured β-titanium alloys' – Journal of Alloys and Compounds.",
            "Artigo de 2019: 'Designing sintering time for a TiSiC compound: a microwave and conventional comparison' – International Journal of Advanced Manufacturing Technology."
        ]
    )

    # --- SLIDE 7: Expertise Complementar ---
    adicionar_slide_topicos(
        "Expertise Complementar",
        [
            "Física e Mecânica: Amplo domínio em cinemática, torque, momento de inércia e hidrostática.",
            "Automação Industrial: Conhecimento prático em sistemas de controle e programação de CLPs, incluindo softwares de simulação."
        ]
    )

    # --- SLIDE 8: Encerramento ---
    adicionar_slide_topicos(
        "Obrigado!",
        [
            "Aberto a perguntas, colaborações em pesquisa e projetos educacionais.",
            "Contato: [Seu Telefone / WhatsApp]",
            "Currículo Lattes: [Link para o Lattes]"
        ]
    )

    # Salva o arquivo final
    nome_arquivo = 'Apresentacao_Lucas_Dias_Calado.pptx'
    prs.save(nome_arquivo)
    print(f"Apresentação '{nome_arquivo}' criada com sucesso!")

# Executa o código
if __name__ == "__main__":
    criar_apresentacao()